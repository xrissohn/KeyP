#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build a vertical (9:16) editable PPTX pitch deck for KeyP using real app screenshots."""
import os
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

BASE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(BASE, "..", "keyp_shots")
ASSETS = os.path.join(BASE, "..", "..", "artifacts", "keyp", "assets", "images")
OUT = os.path.join(BASE, "..", "KAIST OverEdge KeyP \uc18c\uac1c \uc2ac\ub77c\uc774\ub4dc_KeyP.pptx")

# ---- palette ----
BG      = "0A0E1A"
BG2     = "10162A"
CARD    = "161C30"
CARD2   = "1C2540"
INDIGO  = "5B7FFF"
DEEP    = "4F46E5"
PINK    = "FF6B8A"
WHITE   = "FFFFFF"
MUTED   = "9AA4C0"
MUTED2  = "6B7596"
GREEN   = "34D399"
LINE    = "26304C"

FONT = "Malgun Gothic"

# ---- rounded screenshots ----
ROUND_DIR = os.path.join(BASE, "rounded")
os.makedirs(ROUND_DIR, exist_ok=True)

def round_corners(src, dst, radius_frac=0.085):
    im = Image.open(src).convert("RGBA")
    w, h = im.size
    r = int(min(w, h) * radius_frac)
    mask = Image.new("L", (w, h), 0)
    d = ImageDraw.Draw(mask)
    d.rounded_rectangle([0, 0, w, h], radius=r, fill=255)
    im.putalpha(mask)
    im.save(dst)
    return dst

SCREENS = {}
for name in ["01-onboarding", "02-add", "03-match", "04-pricing", "05-interests", "06-profile"]:
    src = os.path.join(SHOTS, name + ".jpg")
    dst = os.path.join(ROUND_DIR, name + ".png")
    round_corners(src, dst)
    SCREENS[name] = dst

# ---- presentation (9:16 vertical) ----
prs = Presentation()
prs.slide_width = Inches(7.5)
prs.slide_height = Inches(13.333)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
M = Inches(0.62)
CW = SW - 2 * M


def rgb(h):
    return RGBColor.from_string(h)


def kfont(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def set_alpha(fill_elem, alpha_pct):
    srgb = fill_elem.find(qn("a:srgbClr"))
    if srgb is not None:
        a = srgb.makeelement(qn("a:alpha"), {})
        a.set("val", str(int(alpha_pct * 1000)))
        srgb.append(a)


def soft_shadow(shape, blur=0.18, dist=0.10, alpha=58, dir=5400000):
    spPr = shape._element.spPr
    eff = spPr.find(qn("a:effectLst"))
    if eff is None:
        eff = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(eff)
    sh = eff.makeelement(qn("a:outerShdw"), {})
    sh.set("blurRad", str(Inches(blur)))
    sh.set("dist", str(Inches(dist)))
    sh.set("dir", str(dir))
    sh.set("rotWithShape", "0")
    clr = sh.makeelement(qn("a:srgbClr"), {})
    clr.set("val", "000000")
    al = clr.makeelement(qn("a:alpha"), {})
    al.set("val", str(alpha * 1000))
    clr.append(al)
    sh.append(clr)
    eff.append(sh)


def bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = rgb(color)
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def glow(slide, cx, cy, size, color, alpha=12):
    e = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - size / 2), int(cy - size / 2), int(size), int(size))
    e.fill.solid()
    e.fill.fore_color.rgb = rgb(color)
    set_alpha(e.fill.fore_color._xFill.find(qn("a:solidFill")) if False else e.fill._xPr.find(qn("a:solidFill")), alpha)
    e.line.fill.background()
    e.shadow.inherit = False
    return e


def rrect(slide, left, top, w, h, fill=CARD, line=None, line_w=1.0, radius=0.06, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(left), int(top), int(w), int(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        soft_shadow(shp)
    return shp


def rect(slide, left, top, w, h, fill, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top), int(w), int(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def R(text, color=None, bold=None, size=None, italic=None):
    return {"t": text, "c": color, "b": bold, "s": size, "i": italic}


def text(slide, left, top, w, h, paras, size=18, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.12, space=6, wrap=True):
    tb = slide.shapes.add_textbox(int(left), int(top), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(paras, str):
        paras = [[R(paras)]]
    elif paras and isinstance(paras[0], dict):
        paras = [paras]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(space)
        p.space_before = Pt(0)
        if isinstance(para, str):
            para = [R(para)]
        for run in para:
            r = p.add_run()
            r.text = run["t"]
            r.font.size = Pt(run["s"] if run["s"] else size)
            r.font.bold = run["b"] if run["b"] is not None else bold
            if run.get("i"):
                r.font.italic = True
            r.font.color.rgb = rgb(run["c"] if run["c"] else color)
            kfont(r)
    return tb


def phone(slide, key, left, top, h):
    img = SCREENS[key]
    iw, ih = Image.open(img).size
    w = int(h * iw / ih)
    pad = Inches(0.07)
    bez = rrect(slide, left - pad, top - pad, w + 2 * pad, h + 2 * pad,
                fill="05070F", line=LINE, line_w=1.25, radius=0.085, shadow=True)
    slide.shapes.add_picture(img, int(left), int(top), int(w), int(h))
    return w


def kicker(slide, top, label, color=INDIGO):
    bar = rect(slide, M, top + Inches(0.02), Inches(0.30), Inches(0.07), color)
    text(slide, M + Inches(0.42), top - Inches(0.05), CW, Inches(0.4),
         [[R(label, color=color, bold=True, size=13)]], align=PP_ALIGN.LEFT)


def page_no(slide, n, label):
    text(slide, M, SH - Inches(0.62), CW, Inches(0.3),
         [[R("KeyP", color=INDIGO, bold=True, size=11), R("  \u00b7  KAIST OverEdge \ucc3d\uc5c5 \uc544\uc774\ub514\uc5b4", color=MUTED2, size=11)]],
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, M, SH - Inches(0.62), CW, Inches(0.3),
         [[R(f"{label}  \u00b7  {n:02d}", color=MUTED2, size=11)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def new_slide(color=BG):
    s = prs.slides.add_slide(BLANK)
    bg(s, color)
    return s


# ============ SLIDE 1 — COVER ============
s = new_slide()
glow(s, SW * 0.85, SH * 0.10, Inches(6.5), DEEP, 16)
glow(s, SW * 0.12, SH * 0.42, Inches(5.0), PINK, 8)
# logo mark
mark = os.path.join(ASSETS, "keyp-icon-mark.png")
s.shapes.add_picture(mark, int(M), int(Inches(0.95)), int(Inches(0.95)), int(Inches(0.95)))
text(s, M + Inches(1.15), Inches(0.98), CW, Inches(0.95),
     [[R("KeyP", color=WHITE, bold=True, size=30)],
      [R("REAL-TIME INTEREST ALERTS \u00b7 AI AGENTS", color=INDIGO, bold=True, size=11)]],
     anchor=MSO_ANCHOR.MIDDLE, line=1.1, space=2)

text(s, M, Inches(2.75), CW, Inches(0.5),
     [[R("2026 KAIST OverEdge (\uc624\uc5e3) \u00b7 \ucc3d\uc5c5 \uc544\uc774\ub514\uc5b4", color=INDIGO, bold=True, size=15)]])
text(s, M, Inches(3.25), CW, Inches(2.2),
     [[R("\uad00\uc2ec\uc0ac\ub97c \ub4f1\ub85d\ub9cc \ud558\uba74,", color=WHITE, bold=True, size=40)],
      [R("AI \uc5d0\uc774\uc804\ud2b8 \ud300\uc774 \ub300\uc2e0", color=WHITE, bold=True, size=40)],
      [R("\ucc3e\uc544 \uc54c\ub824\uc8fc\ub294 ", color=WHITE, bold=True, size=40), R("\uc18d\ubcf4 \uc571", color=INDIGO, bold=True, size=40)]],
     line=1.14, space=2)
text(s, M, Inches(5.55), CW, Inches(1.0),
     [[R("\uc790\uc5f0\uc5b4 \uad00\uc2ec\uc0ac \u2192 4\uac1c AI \uc5d0\uc774\uc804\ud2b8 \ud30c\uc774\ud504\ub77c\uc778 \u2192 \uc2e0\ub8b0\ub3c4 \uac80\uc99d\ub41c \uc54c\ub9bc.", color=MUTED, size=16)],
      [R("\uac70\uae30\uc5d0 ", color=MUTED, size=16), R("opt-in \uc0c1\ud638\ub9e4\uce6d", color=PINK, bold=True, size=16), R("\uc73c\ub85c \uc0ac\ub78c\uae4c\uc9c0 \uc5f0\uacb0\ud558\ub294 \ubaa8\ubc14\uc77c \uc571.", color=MUTED, size=16)]],
     line=1.3, space=4)

# cover phone (onboarding) bottom-right
phone(s, "01-onboarding", SW - Inches(3.55), Inches(6.8), Inches(5.7))
# left chips
chips = [("4\uac1c", "\uc804\ubb38 AI \uc5d0\uc774\uc804\ud2b8"), ("1\uc778", "AI \ud30c\ud2b8\ub108\uc2ed \ucc3d\uc5c5"), ("MVP", "\uc774\ubbf8 \uc791\ub3d9 \uc911")]
cy = Inches(7.4)
for i, (big, sm) in enumerate(chips):
    top = cy + Inches(1.55) * i
    card = rrect(s, M, top, Inches(2.45), Inches(1.25), fill=CARD, line=LINE, line_w=1.0, radius=0.16)
    text(s, M + Inches(0.25), top + Inches(0.18), Inches(2.0), Inches(0.55),
         [[R(big, color=INDIGO, bold=True, size=26)]])
    text(s, M + Inches(0.25), top + Inches(0.78), Inches(2.0), Inches(0.4),
         [[R(sm, color=MUTED, size=13)]])

text(s, M, SH - Inches(0.95), CW, Inches(0.4),
     [[R("\uc2e0\uccad\uc790 \uc131\uba85 / \uc5f0\ub77d\ucc98 : ________________________", color=MUTED2, size=12)]])

# ============ SLIDE 2 — PROBLEM ============
s = new_slide()
glow(s, SW * 0.9, SH * 0.08, Inches(5.5), PINK, 9)
kicker(s, Inches(0.7), "PROBLEM \u00b7 \ud480\uace0\uc790 \ud558\ub294 \ubb38\uc81c")
text(s, M, Inches(1.15), CW, Inches(1.6),
     [[R("\uc815\ubcf4\ub294 \ub118\uce58\ub294\ub370,", color=WHITE, bold=True, size=33)],
      [R("'\ub0b4 \uad00\uc2ec\uc0ac'\uc758 ", color=WHITE, bold=True, size=33), R("\uacb0\uc815\uc801 \uc2e0\ud638", color=PINK, bold=True, size=33)],
      [R("\ub294 \ud56d\uc0c1 \ub2a6\uac8c \ub3c4\ucc29\ud55c\ub2e4.", color=WHITE, bold=True, size=33)]],
     line=1.16, space=2)

probs = [
    ("\uac80\uc0c9\uc740 \uc9c1\uc811 \ucc3e\uc544\uc57c \ud55c\ub2e4", "\ub9e4\ubc88 \uc0ac\ub78c\uc774 \ud0a4\uc6cc\ub4dc\ub97c \ub123\uace0 \ub4a4\uc838\uc57c \ud55c\ub2e4. \uc2e4\uc2dc\uac04\uc131\ub3c4, \ubbf8\ub798 \uc2e0\ud638\ub3c4 \uc5c6\ub2e4."),
    ("\ud0a4\uc6cc\ub4dc \uc54c\ub9bc\uc740 \ub178\uc774\uc988", "\ub2e8\uc21c \uc77c\uce58\ub9cc \ud558\uba74 \uc911\ubcf5\u00b7\uad11\uace0\uac00 \uc30f\ub9ac\uace0, \ud55c\uad6d\uc5b4\u00b7\ub85c\uceec \ucd9c\ucc98\ub294 \uc57d\ud558\ub2e4."),
    ("SNS\ub294 \uad11\uace0\u00b7\uccb4\ub958\uc6a9", "\ucc38\uc5ec\u00b7\uad11\uace0 \uc911\uc2ec \ud53c\ub4dc\ub294 \ud50c\ub7ab\ud3fc \uc774\uc775\uc774 \uc6b0\uc120. \ub0b4 \uad00\uc2ec\uc0ac\ubcf4\ub2e4 \uc54c\uace0\ub9ac\uc998\uc774 \uc6b0\uc120\uc774\ub2e4."),
    ("\ud754\uc5b4\uc9c4 \ucd9c\ucc98, \uac80\uc99d \ubd88\uac00", "\uc5ec\ub7ec \ucd9c\ucc98\ub97c \uac1c\uc778\uc774 \uc77c\uc77c\uc774 \ubaa8\ub2c8\ud130\ub9c1\u00b7\uad50\ucc28\uac80\uc99d\ud558\ub294 \uac83\uc740 \uc2dc\uac04\u00b7\ube44\uc6a9 \ubd88\uac00\ub2a5\ud558\ub2e4."),
]
top = Inches(3.05)
ch = Inches(1.55)
gap = Inches(0.28)
for i, (t, d) in enumerate(probs):
    y = top + (ch + gap) * i
    rrect(s, M, y, CW, ch, fill=CARD, line=LINE, line_w=1.0, radius=0.10)
    b = rrect(s, M + Inches(0.3), y + Inches(0.32), Inches(0.62), Inches(0.62), fill=PINK, radius=0.28)
    text(s, M + Inches(0.3), y + Inches(0.32), Inches(0.62), Inches(0.62),
         [[R(str(i + 1), color=WHITE, bold=True, size=22)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, M + Inches(1.25), y + Inches(0.24), CW - Inches(1.6), Inches(0.5),
         [[R(t, color=WHITE, bold=True, size=18)]])
    text(s, M + Inches(1.25), y + Inches(0.74), CW - Inches(1.6), Inches(0.7),
         [[R(d, color=MUTED, size=13.5)]], line=1.25)

qy = top + (ch + gap) * 4 + Inches(0.05)
qb = rrect(s, M, qy, CW, Inches(1.35), fill="14112A", line=PINK, line_w=1.0, radius=0.10)
text(s, M + Inches(0.4), qy + Inches(0.22), CW - Inches(0.8), Inches(1.0),
     [[R("\u201c\ud754\uc5b4\uc9c4 \uc138\uc0c1\uc758 \uc815\ubcf4 \uc18d\uc5d0\uc11c, \ub0b4 \uad00\uc2ec\uc0ac\uc758 \uacb0\uc815\uc801 \uc2e0\ud638\ub97c", color=WHITE, bold=True, size=17)],
      [R("\uac00\uc7a5 \uba3c\uc800\u00b7\uc815\ud655\ud558\uac8c \ubc1b\uc744 \ubc29\ubc95\uc774 \uc5c6\ub2e4.\u201d", color=WHITE, bold=True, size=17)]],
     line=1.3, space=2)
page_no(s, 2, "Problem")

# ============ SLIDE 3 — SOLUTION ============
s = new_slide()
glow(s, SW * 0.12, SH * 0.12, Inches(5.5), INDIGO, 12)
kicker(s, Inches(0.7), "SOLUTION \u00b7 \ud574\uacb0\ucc45")
text(s, M, Inches(1.15), CW, Inches(1.2),
     [[R("\ub4f1\ub85d\ub9cc \ud558\uba74,", color=WHITE, bold=True, size=33)],
      [R("AI\uac00 ", color=WHITE, bold=True, size=33), R("\ub300\uc2e0 \ucc3e\uc544 \uc54c\ub824\uc900\ub2e4", color=INDIGO, bold=True, size=33)]],
     line=1.16, space=2)
phone(s, "01-onboarding", M + Inches(0.15), Inches(3.0), Inches(5.7))

bx = M + Inches(3.55)
bw = SW - bx - M
feats = [
    ("\uc790\uc5f0\uc5b4\ub85c \uad00\uc2ec\uc0ac \ub4f1\ub85d", "\u201c\uba87 \ud0a4\uc6cc\ub4dc\u201d \uac00 \uc544\ub2c8\ub77c \uc77c\uc0c1 \ubb38\uc7a5\uc73c\ub85c \uc801\uc73c\uba74 \ub41c\ub2e4."),
    ("\uad6c\uc870\ud654 \u00b7 \uc804\ub7b5 \uc218\ub9bd", "AI\uac00 \uc758\ub3c4\ub97c \uc77d\uace0 '\uc5b4\ub514\uc11c \uc5b4\ub5bb\uac8c \ucc3e\uc744\uc9c0' \uacc4\ud68d\uc744 \uc138\uc6b4\ub2e4."),
    ("\uc2e0\ub8b0\ub3c4 \uac80\uc99d\ub41c \uc54c\ub9bc", "\ucd9c\ucc98 \uad50\ucc28\uac80\uc99d \ud6c4 \uc810\uc218\uac00 \ub192\uc740 \uac83\ub9cc \uc694\uc57d\ud574 \uc804\ub2ec\ud55c\ub2e4."),
    ("\uadf8\ub300\ub85c \uc18d\ubcf4 \ud53c\ub4dc", "\uc0c8 \uc18c\uc2dd\uc774 \uc624\uba74 \ud478\uc2dc \uc54c\ub9bc. \uc9c1\uc811 \ucc3e\uc9c0 \uc54a\uc544\ub3c4 \ub41c\ub2e4."),
]
text(s, bx, Inches(2.95), bw, Inches(0.5),
     [[R("\ud575\uc2ec \uacbd\ud5d8", color=INDIGO, bold=True, size=15)]])
yy = Inches(3.5)
for i, (t, d) in enumerate(feats):
    y = yy + Inches(1.5) * i
    rrect(s, bx, y, bw, Inches(1.32), fill=CARD, line=LINE, line_w=1.0, radius=0.12)
    text(s, bx + Inches(0.22), y + Inches(0.18), bw - Inches(0.4), Inches(0.45),
         [[R(f"{i+1}. ", color=INDIGO, bold=True, size=16), R(t, color=WHITE, bold=True, size=16)]])
    text(s, bx + Inches(0.22), y + Inches(0.66), bw - Inches(0.4), Inches(0.6),
         [[R(d, color=MUTED, size=12.5)]], line=1.25)
page_no(s, 3, "Solution")

# ============ SLIDE 4 — HOW IT WORKS (AGENTS) ============
s = new_slide()
glow(s, SW * 0.9, SH * 0.9, Inches(6.0), DEEP, 14)
kicker(s, Inches(0.7), "HOW IT WORKS \u00b7 4\uac1c AI \uc5d0\uc774\uc804\ud2b8")
text(s, M, Inches(1.15), CW, Inches(1.2),
     [[R("'\uac80\uc0c9'\uc774 \uc544\ub2c8\ub77c ", color=WHITE, bold=True, size=31), R("'\ucd94\ub860 \ud6c4 \uc88c\ud600 \ucc3e\uae30'", color=INDIGO, bold=True, size=31)]],
     line=1.15)

agents = [
    ("Planner", "GPT", "\uad00\uc2ec\uc0ac \uad6c\uc870\ud654\u00b7\uac80\uc0c9 \uc804\ub7b5 \uc218\ub9bd", INDIGO),
    ("Collector", "Perplexity", "\uc804\ub7b5\ub300\ub85c \uc720\ub9dd \ucd9c\ucc98\ub9cc \uc2e4\uc81c \uc218\uc9d1", DEEP),
    ("Verifier", "Claude", "\ucd9c\ucc98 \uad50\ucc28\uac80\uc99d \u00b7 \uc2e0\ub8b0\ub3c4 \uc810\uc218", "7C6BFF"),
    ("Deliverer", "Sort\u00b7Dedup", "\ucd5c\uc2e0\uc131\u00b7\uc2e0\ub8b0\ub3c4 \uc815\ub82c \u00b7 \uc911\ubcf5\uc81c\uac70", PINK),
]
top = Inches(2.55)
ah = Inches(1.25)
agap = Inches(0.55)
for i, (nm, sub, d, col) in enumerate(agents):
    y = top + (ah + agap) * i
    rrect(s, M, y, CW, ah, fill=CARD, line=col, line_w=1.25, radius=0.12)
    badge = rrect(s, M + Inches(0.28), y + Inches(0.30), Inches(0.66), Inches(0.66), fill=col, radius=0.5)
    text(s, M + Inches(0.28), y + Inches(0.30), Inches(0.66), Inches(0.66),
         [[R(str(i + 1), color=WHITE, bold=True, size=22)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, M + Inches(1.2), y + Inches(0.22), CW - Inches(1.5), Inches(0.5),
         [[R(nm, color=WHITE, bold=True, size=20), R(f"   {sub}", color=col, bold=True, size=13)]])
    text(s, M + Inches(1.2), y + Inches(0.72), CW - Inches(1.5), Inches(0.4),
         [[R(d, color=MUTED, size=13.5)]])
    if i < 3:
        ar = text(s, M, y + ah - Inches(0.02), CW, agap,
                  [[R("\u2193", color=col, bold=True, size=20)]], align=PP_ALIGN.CENTER)

phy = top + (ah + agap) * 4 - Inches(0.15)
note = rrect(s, M, phy, CW, Inches(1.25), fill="101A14", line=GREEN, line_w=1.0, radius=0.10)
text(s, M + Inches(0.35), phy + Inches(0.2), CW - Inches(0.7), Inches(0.95),
     [[R("\ud575\uc2ec \ucca0\ud559 \u2014 AlphaGo\uc2dd \ud6c4\ubcf4 \ucd95\uc18c", color=GREEN, bold=True, size=15)],
      [R("\ubaa8\ub4e0 \ud0a4\uc6cc\ub4dc\ub97c \ubaa8\ub4e0 \uacf3\uc5d0 \uacc4\uc18d \uae4b\uc9c0 \uc54a\ub294\ub2e4. \uac00\ub2a5\uc131 \ub192\uc740 \ud6c4\ubcf4\ubd80\ud130 \uc88c\ud600 \u2192 \ube44\uc6a9\u2193 \uc18d\ub3c4\u00b7\uc815\ud655\ub3c4\u2191", color=MUTED, size=12.5)]],
     line=1.3, space=3)
page_no(s, 4, "How it works")

# ============ SLIDE 5 — PRODUCT (MVP) ============
s = new_slide()
glow(s, SW * 0.5, SH * 0.06, Inches(6.5), INDIGO, 10)
kicker(s, Inches(0.7), "PRODUCT \u00b7 \uc774\ubbf8 \uc791\ub3d9\ud558\ub294 \uc81c\ud488")
text(s, M, Inches(1.15), CW, Inches(1.2),
     [[R("\uae30\ud68d\uc548\uc774 \uc544\ub2c8\ub77c, ", color=WHITE, bold=True, size=31), R("\uc791\ub3d9\ud558\ub294 MVP", color=INDIGO, bold=True, size=31)]],
     line=1.15)
# two phones row
ph_h = Inches(5.6)
phone(s, "05-interests", M + Inches(0.2), Inches(2.5), ph_h)
phone(s, "02-add", SW / 2 + Inches(0.35), Inches(2.5), ph_h)
text(s, M, Inches(2.5) + ph_h + Inches(0.18), CW / 2, Inches(0.4),
     [[R("\u2460 \uad00\uc2ec\uc0ac \uad00\ub9ac \u00b7 \uc790\ub3d9 \uc218\uc9d1", color=MUTED, size=12.5)]], align=PP_ALIGN.CENTER)
text(s, SW / 2, Inches(2.5) + ph_h + Inches(0.18), CW / 2, Inches(0.4),
     [[R("\u2461 AI \uc5d0\uc774\uc804\ud2b8 \uad00\uc2ec\uc0ac \ubd84\uc11d", color=MUTED, size=12.5)]], align=PP_ALIGN.CENTER)

dy = Inches(9.05)
done = ["4\uac1c AI \uc5d0\uc774\uc804\ud2b8 \uc2e4\uc2dc\uac04 \ud30c\uc774\ud504\ub77c\uc778", "\uc2e4\uc2dc\uac04 \ud478\uc2dc \uc54c\ub9bc (PWA Web Push)",
        "\uc758\ubbf8 \uae30\ubc18 \uc911\ubcf5\uc81c\uac70 \u00b7 \uc8fd\uc740 \ub9c1\ud06c \ucc28\ub2e8", "\uc88b\uc544\uc694/\ubcc4\ub85c\uc608\uc694 \ud53c\ub4dc\ubc31 \ud559\uc2b5",
        "\ud55c\uad6d\uc5b4 UI + \uae00\ub85c\ubc8c \uac80\uc0c9\u00b7\ubc88\uc5ed", "\ube44\uc6a9 \ucd5c\uc801\ud654 (\uce90\uc2f1\u00b7\ucf54\uc5bc\ub808\uc2f1)"]
cols = 2
cwid = CW / 2 - Inches(0.1)
for i, d in enumerate(done):
    rr = i // cols
    cc = i % cols
    x = M + cc * (cwid + Inches(0.2))
    y = dy + rr * Inches(0.7)
    rrect(s, x, y, cwid, Inches(0.58), fill=CARD, line=LINE, line_w=0.75, radius=0.18)
    text(s, x + Inches(0.18), y, cwid - Inches(0.3), Inches(0.58),
         [[R("\u2713  ", color=GREEN, bold=True, size=13), R(d, color=WHITE, size=12.5)]],
         anchor=MSO_ANCHOR.MIDDLE)
page_no(s, 5, "Product")

# ============ SLIDE 6 — MATCHING ============
s = new_slide()
glow(s, SW * 0.85, SH * 0.15, Inches(5.5), PINK, 12)
kicker(s, Inches(0.7), "MATCHING \u00b7 \uc0c1\ud638\ub9e4\uce6d", color=PINK)
text(s, M, Inches(1.15), CW, Inches(1.2),
     [[R("\uac19\uc740 \uad00\uc2ec\uc0ac\ub97c \uac00\uc9c4 \uc0ac\ub78c\uacfc", color=WHITE, bold=True, size=31)],
      [R("\uc591\ubc29\ud5a5 \ub3d9\uc758 \uc2dc\uc5d0\ub9cc ", color=WHITE, bold=True, size=31), R("\uc548\uc804\ud558\uac8c \uc5f0\uacb0", color=PINK, bold=True, size=31)]],
     line=1.16, space=2)
phone(s, "03-match", M + Inches(0.15), Inches(3.0), Inches(5.7))
bx = M + Inches(3.55)
bw = SW - bx - M
mfeats = [
    ("opt-in \uc591\ubc29\ud5a5", "\uc11c\ub85c \ub3d9\uc758\ud55c \uacbd\uc6b0\uc5d0\ub9cc \uc5f0\uacb0. \uc6d0\uce58 \uc54a\uc73c\uba74 \uc5f0\uacb0\ub418\uc9c0 \uc54a\ub294\ub2e4."),
    ("\uc548\uc804 \uc7a5\uce58", "\uc2e0\uace0\u00b7\ucc28\ub2e8 \uc9c0\uc6d0. \uad00\uc2ec\uc0ac \uae30\ubc18\uc774\ub77c \ubd80\ub2f4\uc774 \uc801\ub2e4."),
    ("\ub124\ud2b8\uc6cc\ud06c \ud6a8\uacfc", "\uc0ac\uc6a9\uc790\uac00 \ub298\uc218\ub85d \ub9e4\uce6d \uac00\uce58\uac00 \ucee4\uc9c0\ub294 \uad6c\uc870."),
    ("\uc54c\ub9bc + \uc5f0\uacb0", "'\uc815\ubcf4'\uc640 '\uc0ac\ub78c'\uc744 \ud55c \uc571\uc5d0\uc11c \ubaa8\ub450 \uc81c\uacf5."),
]
text(s, bx, Inches(2.95), bw, Inches(0.5),
     [[R("\uc65c \uac15\ub825\ud55c\uac00", color=PINK, bold=True, size=15)]])
for i, (t, d) in enumerate(mfeats):
    y = Inches(3.5) + Inches(1.5) * i
    rrect(s, bx, y, bw, Inches(1.32), fill=CARD, line=LINE, line_w=1.0, radius=0.12)
    text(s, bx + Inches(0.22), y + Inches(0.18), bw - Inches(0.4), Inches(0.45),
         [[R(t, color=WHITE, bold=True, size=16)]])
    text(s, bx + Inches(0.22), y + Inches(0.64), bw - Inches(0.4), Inches(0.6),
         [[R(d, color=MUTED, size=12.5)]], line=1.25)
page_no(s, 6, "Matching")

# ============ SLIDE 7 — DIFFERENTIATION ============
s = new_slide()
glow(s, SW * 0.15, SH * 0.9, Inches(6.0), DEEP, 12)
kicker(s, Inches(0.7), "EDGE \u00b7 \uc81c\ud488 \ucc28\ubcc4\uc131")
text(s, M, Inches(1.15), CW, Inches(1.2),
     [[R("4\uac00\uc9c0 ", color=WHITE, bold=True, size=33), R("'Edge'", color=INDIGO, bold=True, size=33)]],
     line=1.15)
edges = [
    ("A", "\uba40\ud2f0 \uc5d0\uc774\uc804\ud2b8 \uc815\ubc00\ub3c4", "\uc218\uc9d1\u00b7\uac80\uc99d\u00b7\uc804\ub2ec\uc744 \ubd84\uc5c5\ud574 \ub178\uc774\uc988\ub97c \uac78\ub7ec\ub0b8\ub2e4. \ub2e8\uc21c \ud0a4\uc6cc\ub4dc\uac00 \uc544\ub2cc \ub9e5\ub77d\u00b7\uc758\ub3c4 \uae30\ubc18.", INDIGO),
    ("B", "\uac80\uc99d\ub41c \uc2e0\ub8b0\ub3c4 + \ucd9c\ucc98", "\ubaa8\ub4e0 \uc54c\ub9bc\uc5d0 \uc2e0\ub8b0\ub3c4 \uc810\uc218\u00b7\uadfc\uac70\u00b7\uc6d0\ubb38 \ub9c1\ud06c. '\ubbff\uace0 \ub204\ub974\ub294' \uacbd\ud5d8.", "7C6BFF"),
    ("C", "\ud55c\uad6d\uc5b4 \uc6b0\uc120 \u00b7 \uae00\ub85c\ubc8c", "\ud55c\uad6d\uc5b4 UI\ub85c \uc4f0\uba74\uc11c \uc804 \uc138\uacc4 \ucd9c\ucc98\ub97c \uac80\uc0c9\u00b7\uc790\ub3d9 \ubc88\uc5ed\ud55c\ub2e4.", GREEN),
    ("D", "opt-in \uc0c1\ud638\ub9e4\uce6d", "\uac19\uc740 \uad00\uc2ec\uc0ac\ub97c \uc591\ubc29\ud5a5 \ub3d9\uc758 \uc2dc\uc5d0\ub9cc \uc5f0\uacb0. \ub124\ud2b8\uc6cc\ud06c \ud6a8\uacfc.", PINK),
]
top = Inches(2.5)
ch2 = Inches(2.15)
gap2 = Inches(0.3)
cwid = CW / 2 - Inches(0.15)
for i, (lab, t, d, col) in enumerate(edges):
    rr = i // 2
    cc = i % 2
    x = M + cc * (cwid + Inches(0.3))
    y = top + rr * (ch2 + gap2)
    rrect(s, x, y, cwid, ch2, fill=CARD, line=LINE, line_w=1.0, radius=0.10)
    b = rrect(s, x + Inches(0.25), y + Inches(0.25), Inches(0.62), Inches(0.62), fill=col, radius=0.28)
    text(s, x + Inches(0.25), y + Inches(0.25), Inches(0.62), Inches(0.62),
         [[R(lab, color=WHITE, bold=True, size=22)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.25), y + Inches(1.02), cwid - Inches(0.5), Inches(0.5),
         [[R(t, color=WHITE, bold=True, size=16.5)]], line=1.1)
    text(s, x + Inches(0.25), y + Inches(1.45), cwid - Inches(0.5), Inches(0.65),
         [[R(d, color=MUTED, size=12)]], line=1.22)

cy = top + 2 * (ch2 + gap2) + Inches(0.05)
rrect(s, M, cy, CW, Inches(1.4), fill="0E1530", line=INDIGO, line_w=1.0, radius=0.10)
text(s, M + Inches(0.35), cy + Inches(0.24), CW - Inches(0.7), Inches(1.0),
     [[R("\uacbd\uc7c1\ub825 \uc694\uc57d", color=INDIGO, bold=True, size=14)],
      [R("KeyP\ub294 ", color=WHITE, size=15), R("'\uc544\uc774\ub514\uc5b4'\uac00 \uc544\ub2c8\ub77c \uc774\ubbf8 \uad6c\ud604\ub418\uc5b4 \uc791\ub3d9\ud558\ub294 \uc81c\ud488", color=WHITE, bold=True, size=15),
       R("\uc774\ub2e4. \uc774\uac83\uc774 \ud575\uc2ec \uacbd\uc7c1\ub825\uc774\ub2e4.", color=WHITE, size=15)]],
     line=1.3, space=4)
page_no(s, 7, "Edge")

# ============ SLIDE 8 — BUSINESS MODEL ============
s = new_slide()
glow(s, SW * 0.9, SH * 0.1, Inches(5.5), INDIGO, 12)
kicker(s, Inches(0.7), "BUSINESS \u00b7 \uc218\uc775 \ubaa8\ub378")
text(s, M, Inches(1.15), CW, Inches(1.2),
     [[R("\uad6c\ub3c5 4\ub2e8\uacc4 \u2014 ", color=WHITE, bold=True, size=31), R("Free \u2192 Power", color=INDIGO, bold=True, size=31)]],
     line=1.15)
phone(s, "04-pricing", M + Inches(0.15), Inches(2.95), Inches(5.7))
bx = M + Inches(3.55)
bw = SW - bx - M
plans = [
    ("Free", "\uad00\uc2ec\uc0ac 1 \u00b7 1\uc2dc\uac04", MUTED),
    ("Basic", "\uad00\uc2ec\uc0ac 5 \u00b7 15\ubd84", INDIGO),
    ("Pro", "\uace0\ube48\ub3c4 \u00b7 \uc18d\ubcf4 \uc54c\ub9bc", "7C6BFF"),
    ("Power", "\ucd5c\ub300 \u00b7 \ucd5c\uace0\ube48\ub3c4+\ubd80\uc2a4\ud2b8", PINK),
]
text(s, bx, Inches(2.92), bw, Inches(0.5),
     [[R("\uad00\uc2ec\uc0ac \uc218\u00b7\uc54c\ub9bc \ube48\ub3c4\ub85c \ucc28\ub4f1", color=INDIGO, bold=True, size=14)]])
for i, (nm, d, col) in enumerate(plans):
    y = Inches(3.5) + Inches(1.5) * i
    rrect(s, bx, y, bw, Inches(1.32), fill=CARD, line=col, line_w=1.1, radius=0.12)
    text(s, bx + Inches(0.22), y + Inches(0.2), bw - Inches(0.4), Inches(0.5),
         [[R(nm, color=col, bold=True, size=19)]])
    text(s, bx + Inches(0.22), y + Inches(0.7), bw - Inches(0.4), Inches(0.5),
         [[R(d, color=MUTED, size=13)]])
text(s, M, Inches(9.6), CW, Inches(0.5),
     [[R("\uc5c5\uc12c\ub110\u00b7\uc218\uc0c1 11M KRW + \uc2dc\ub4dc \ud22c\uc790 \ucd5c\ub300 100M KRW \uc5f0\uacc4 \uac00\ub2a5 (OverEdge \ud2b8\ub799)", color=MUTED2, size=12)]])
page_no(s, 8, "Business")

# ============ SLIDE 9 — 1\uc778 \ucc3d\uc5c5 x AI ============
s = new_slide()
glow(s, SW * 0.5, SH * 0.12, Inches(6.5), DEEP, 14)
kicker(s, Inches(0.7), "WHY KeyP \u00b7 1\uc778 \ucc3d\uc5c5 \u00d7 AI")
text(s, M, Inches(1.2), CW, Inches(2.0),
     [[R("KeyP \uc790\uccb4\uac00", color=WHITE, bold=True, size=34)],
      [R("'AI\ub97c \uacf5\ub3d9\ucc3d\uc5c5\ud300\uc73c\ub85c \uc0bc\uc740", color=INDIGO, bold=True, size=34)],
      [R("1\uc778 \ucc3d\uc5c5'", color=INDIGO, bold=True, size=34), R(" \uc758 \uc0ac\ub840\ub2e4.", color=WHITE, bold=True, size=34)]],
     line=1.16, space=2)
text(s, M, Inches(3.55), CW, Inches(1.0),
     [[R("\ucc3d\uc5c5\uc790\ub294 \ubcc4\ub3c4 \ud300 \uc5c6\uc774 AI \uc5d0\uc774\uc804\ud2b8\ub97c \uc2e4\uc9c8\uc801 \ub3d9\ub8cc\ub85c \uc0bc\uc544 \ubaa8\ubc14\uc77c \uc571\u00b7\ubc31\uc5d4\ub4dc\u00b7AI \ud30c\uc774\ud504\ub77c\uc778 \uc804\uccb4\ub97c \ud63c\uc790 \uad6c\ucd95\ud588\ub2e4.", color=MUTED, size=15)]],
     line=1.35)

fit = [
    ("\ubb38\uc81c\uc778\uc2dd", "'\uc815\ubcf4 \uacfc\uc789 \uc18d \ubc1c\uacac \ubd80\uc871'\uc744 \uad6c\uc870\uc801\uc73c\ub85c \uc815\uc758"),
    ("\ucc3d\uc5c5\uc758\uc9c0\u00b7\uc2e4\ud604\uac00\ub2a5", "\uae30\ud68d\uc774 \uc544\ub2cc \uc791\ub3d9\ud558\ub294 MVP \u00b7 \uad6c\ub3c5 BM \uc644\uc131"),
    ("AI \ud65c\uc6a9 \uc5ed\ub7c9", "AI\ub97c \uc81c\ud488 \uc5d4\uc9c4\uc774\uc790 \uacf5\ub3d9\ucc3d\uc5c5\ud300\uc73c\ub85c \ub3d9\uc2dc \ud65c\uc6a9"),
    ("AX \uad6c\ud604 (2\ucc28)", "\uba40\ud2f0 \uc5d0\uc774\uc804\ud2b8 \uc624\ucf00\uc2a4\ud2b8\ub808\uc774\uc158\uc744 \uc2e4\uc81c \ucf54\ub4dc\ub85c \uad6c\ud604"),
    ("\uae30\uc220\u00b7\uc2dc\uc7a5 (2\ucc28)", "\ube44\uc6a9 \ucd5c\uc801\ud654 + \ub124\ud2b8\uc6cc\ud06c \ud6a8\uacfc\ud615 \ud655\uc7a5 \uacbd\ub85c"),
]
top = Inches(4.85)
for i, (k, v) in enumerate(fit):
    y = top + Inches(1.05) * i
    rrect(s, M, y, CW, Inches(0.9), fill=CARD, line=LINE, line_w=1.0, radius=0.14)
    rect(s, M, y + Inches(0.18), Inches(0.09), Inches(0.54), INDIGO)
    text(s, M + Inches(0.35), y, Inches(2.55), Inches(0.9),
         [[R(k, color=WHITE, bold=True, size=15)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, M + Inches(2.95), y, CW - Inches(3.2), Inches(0.9),
         [[R(v, color=MUTED, size=13)]], anchor=MSO_ANCHOR.MIDDLE, line=1.15)
page_no(s, 9, "Why KeyP")

# ============ SLIDE 10 — CLOSING ============
s = new_slide()
glow(s, SW * 0.5, SH * 0.4, Inches(8.0), DEEP, 16)
glow(s, SW * 0.5, SH * 0.85, Inches(5.5), PINK, 9)
s.shapes.add_picture(mark, int(SW / 2 - Inches(0.65)), int(Inches(2.4)), int(Inches(1.3)), int(Inches(1.3)))
text(s, M, Inches(4.0), CW, Inches(0.6),
     [[R("KeyP", color=WHITE, bold=True, size=34)]], align=PP_ALIGN.CENTER)
text(s, M, Inches(4.8), CW, Inches(2.4),
     [[R("\ub0b4 \uad00\uc2ec\uc0ac\ub97c", color=WHITE, bold=True, size=34)],
      [R("\uac00\uc7a5 \uba3c\uc800\u00b7\uc815\ud655\ud558\uac8c", color=INDIGO, bold=True, size=34)],
      [R("\uc54c\ub824\uc8fc\ub294 \uae00\ub85c\ubc8c \uc11c\ube44\uc2a4\ub85c", color=WHITE, bold=True, size=34)]],
     align=PP_ALIGN.CENTER, line=1.2, space=2)
text(s, M, Inches(7.6), CW, Inches(1.0),
     [[R("\uc774\ubbf8 \uc791\ub3d9\ud558\ub294 AI \uc5d0\uc774\uc804\ud2b8 \uc81c\ud488 + \uacf5\ub3d9\ucc3d\uc5c5\ud300\uc73c\ub85c\uc11c\uc758 AI.", color=MUTED, size=16)],
      [R("OverEdge\uc758 \uae30\uc220\u00b7\ub124\ud2b8\uc6cc\ud06c\u00b7\ud22c\uc790\uc640 \ub9cc\ub098\uba74 \uae00\ub85c\ubc8c\ub85c \uc131\uc7a5\ud55c\ub2e4.", color=MUTED, size=16)]],
     align=PP_ALIGN.CENTER, line=1.35, space=3)
rrect(s, M + Inches(1.0), Inches(9.4), CW - Inches(2.0), Inches(0.95), fill=CARD, line=INDIGO, line_w=1.0, radius=0.3)
text(s, M + Inches(1.0), Inches(9.4), CW - Inches(2.0), Inches(0.95),
     [[R("2026 KAIST OverEdge \u00b7 \ucc3d\uc5c5 \uc544\uc774\ub514\uc5b4 \uae30\uc220\uc11c", color=WHITE, bold=True, size=15)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, M, Inches(10.7), CW, Inches(0.4),
     [[R("\uc2e0\uccad\uc790 \uc131\uba85 / \uc5f0\ub77d\ucc98 : ________________________", color=MUTED2, size=12)]], align=PP_ALIGN.CENTER)

prs.save(OUT)
print("SAVED:", OUT)
print("slides:", len(prs.slides._sldIdLst))
